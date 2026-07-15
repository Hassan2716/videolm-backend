"""
AI Presentation Generator — creates PPTX from video content.

Bugs fixed:
- Non-modern themes only generating 2 slides: heuristic_plan now cycles
  through content properly when summary is short, using transcript as fallback
- Sentences now always complete

Install python-pptx:
  C:\\Users\\LENOVO\\AppData\\Local\\Programs\\Python\\Python311\\python.exe -m pip install python-pptx
"""
import os, re
from typing import List, Dict, Optional
from loguru import logger
from app.core.config import settings

THEMES = {
    "modern":    {"bg": (15,23,42),     "title": (96,165,250),   "body": (226,232,240), "accent": (139,92,246)},
    "academic":  {"bg": (255,255,255),  "title": (30,58,138),    "body": (31,41,55),    "accent": (37,99,235)},
    "corporate": {"bg": (248,250,252),  "title": (17,24,39),     "body": (55,65,81),    "accent": (59,130,246)},
    "dark":      {"bg": (9,9,11),       "title": (167,139,250),  "body": (212,212,216), "accent": (124,58,237)},
    "light":     {"bg": (255,255,255),  "title": (17,24,39),     "body": (75,85,99),    "accent": (16,185,129)},
    "minimal":   {"bg": (250,250,250),  "title": (0,0,0),        "body": (75,85,99),    "accent": (0,0,0)},
    "ai":        {"bg": (2,6,23),       "title": (34,211,238),   "body": (203,213,225), "accent": (168,85,247)},
}


def _check_pptx():
    try:
        import pptx
        return True
    except ImportError:
        import sys
        raise ImportError(
            f"python-pptx not installed.\n"
            f"Run: {sys.executable} -m pip install python-pptx\n"
            f"Then restart the backend."
        )


def _get_all_sentences(summary: str, transcript: str, min_words: int = 8) -> List[str]:
    """
    Extract clean complete sentences from summary + transcript combined.
    Always returns enough sentences to fill any slide count.
    """
    # Use both summary and transcript for content richness
    combined = ""
    if summary and len(summary.strip()) > 50:
        combined = summary + " " + transcript[:4000]
    else:
        combined = transcript[:6000]

    # Clean markdown
    combined = re.sub(r'#+\s*', '', combined)
    combined = re.sub(r'\*+', '', combined)
    combined = re.sub(r'•\s*', '', combined)
    combined = re.sub(r'\s+', ' ', combined).strip()

    # Split into sentences
    raw = re.split(r'(?<=[.!?])\s+', combined)
    sentences = []
    for s in raw:
        s = s.strip()
        if len(s.split()) < min_words:
            continue
        # Skip filler
        if re.search(r'\b(subscribe|like and|hit the bell|my name is|welcome back)\b', s, re.I):
            continue
        # Ensure complete
        if s and s[-1] not in '.!?':
            s = s + '.'
        sentences.append(s)

    return sentences


def _extract_topics(sentences: List[str], num_topics: int) -> List[Dict]:
    """
    Divide sentences into num_topics groups intelligently.
    Always produces exactly num_topics groups even if content is short.
    """
    if not sentences:
        return [{"title": f"Topic {i+1}", "bullets": ["Content not available."]} for i in range(num_topics)]

    # Repeat sentences if we don't have enough (rare, handles very short videos)
    while len(sentences) < num_topics * 3:
        sentences = sentences + sentences

    # Group into chunks
    chunk_size = max(1, len(sentences) // num_topics)
    topics = []
    for i in range(num_topics):
        start = i * chunk_size
        chunk = sentences[start:start + chunk_size + 2]
        if not chunk:
            chunk = sentences[:3]

        # Title from first sentence (5 words)
        first_words = chunk[0].split()[:6]
        title = " ".join(first_words).rstrip(".,;:!?")
        if len(title) > 55:
            title = title[:52] + "…"

        topics.append({
            "title": f"Topic {i+1}: {title}",
            "bullets": chunk[:3],
            "note": " ".join(chunk[:2]),
        })
    return topics


def _plan_slides(transcript: str, summary: str, num_slides: int) -> List[Dict]:
    """
    Plan exactly num_slides slides.
    Uses transcript + summary to ensure enough content regardless of theme.
    """
    sentences = _get_all_sentences(summary, transcript)
    logger.info(f"Planning {num_slides} slides from {len(sentences)} sentences")

    slides = []

    # Slide 1: Title / Introduction (first 3 sentences)
    intro_bullets = sentences[:3] if len(sentences) >= 3 else sentences
    slides.append({
        "title": "Introduction",
        "bullets": intro_bullets,
        "note": "Opening overview of key topics covered in this video.",
    })

    # Middle slides: distribute remaining content
    remaining = sentences[3:]
    num_content_slides = num_slides - 2  # subtract intro + conclusion
    topics = _extract_topics(remaining, num_content_slides)
    slides.extend(topics)

    # Final slide: Key Takeaways (last 3 sentences)
    takeaway_sents = sentences[-3:] if len(sentences) >= 3 else sentences
    slides.append({
        "title": "Key Takeaways & Conclusions",
        "bullets": takeaway_sents,
        "note": "Summary of the most important concepts from this video.",
    })

    logger.info(f"Planned {len(slides)} slides")
    return slides[:num_slides]


def _match_frame(title: str, bullets: List[str], frames: List[Dict]) -> Optional[str]:
    if not frames:
        return None
    combined = (title + " " + " ".join(bullets)).lower()
    query_words = set(re.findall(r'\b\w{4,}\b', combined))
    best_score, best_path = 0, None
    for frame in frames:
        frame_text = ((frame.get("caption") or "") + " " + (frame.get("ocr_text") or "")).lower()
        frame_words = set(re.findall(r'\b\w{4,}\b', frame_text))
        score = len(query_words & frame_words)
        if score > best_score:
            best_score = score
            best_path = frame.get("frame_path")
    return best_path if best_score > 0 else None


def generate_pptx(
    project_id: str,
    transcript: str,
    summary: str,
    frames: List[Dict],
    num_slides: int = 10,
    theme_name: str = "modern",
    include_images: bool = True,
) -> str:
    """Generate complete PPTX. Returns output path."""
    _check_pptx()

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    theme = THEMES.get(theme_name, THEMES["modern"])

    def rgb(t): return RGBColor(*t)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    # Plan ALL slides upfront using both transcript and summary
    slide_plans = _plan_slides(transcript, summary, num_slides)
    logger.info(f"Generating {len(slide_plans)} slides with theme={theme_name}")

    for slide_idx, plan in enumerate(slide_plans):
        slide = prs.slides.add_slide(BLANK)

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(theme["bg"])

        # Left accent bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(theme["accent"])
        bar.line.fill.background()

        # Slide number
        sn_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
        sn_tf = sn_box.text_frame
        sn_p = sn_tf.paragraphs[0]
        sn_run = sn_p.add_run()
        sn_run.text = f"{slide_idx + 1} / {len(slide_plans)}"
        sn_run.font.size = Pt(9)
        sn_run.font.color.rgb = rgb(theme["body"])

        # Title
        title_text = plan.get("title", f"Slide {slide_idx + 1}")
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(8.5), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = rgb(theme["title"])

        # Separator line
        sep = slide.shapes.add_shape(1, Inches(0.3), Inches(1.35), Inches(12.7), Pt(1.5))
        sep.fill.solid()
        sep.fill.fore_color.rgb = rgb(theme["accent"])
        sep.line.fill.background()

        # Bullets
        bullets = plan.get("bullets", [])
        has_image_space = include_images and frames
        bullet_width = Inches(7.8) if has_image_space else Inches(12.5)

        if bullets:
            btb = slide.shapes.add_textbox(Inches(0.3), Inches(1.55), bullet_width, Inches(5.6))
            btf = btb.text_frame
            btf.word_wrap = True
            for i, bullet_text in enumerate(bullets[:4]):
                p2 = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                run2 = p2.add_run()
                if bullet_text and bullet_text[-1] not in '.!?':
                    bullet_text = bullet_text + '.'
                run2.text = f"• {bullet_text}"
                run2.font.size = Pt(18)
                run2.font.color.rgb = rgb(theme["body"])

        # Frame image
        if has_image_space:
            frame_path = _match_frame(title_text, bullets, frames)
            if frame_path:
                full_path = frame_path if os.path.isabs(frame_path) else \
                            os.path.join(settings.output_dir, frame_path)
                if not os.path.exists(full_path):
                    full_path = os.path.join(settings.output_dir, project_id,
                                             "detected", os.path.basename(frame_path))
                if os.path.exists(full_path):
                    try:
                        slide.shapes.add_picture(
                            full_path, Inches(8.4), Inches(1.55), Inches(4.7), Inches(4.8)
                        )
                    except Exception as img_err:
                        logger.debug(f"Image insert failed: {img_err}")

        # Speaker notes
        note = plan.get("note", "")
        if note:
            slide.notes_slide.notes_text_frame.text = note

    out_dir = os.path.join(settings.export_dir, project_id)
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"presentation_{theme_name}.pptx")
    prs.save(out_path)
    logger.info(f"PPTX saved: {out_path} ({len(slide_plans)} slides)")
    return out_path
